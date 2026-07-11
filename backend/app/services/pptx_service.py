"""PowerPoint <-> PDF conversion service.

- powerpoint_to_pdf: uses LibreOffice headless
- pdf_to_powerpoint: renders each PDF page as a high-quality image slide
  via pdf2image + python-pptx.  Guarantees visual fidelity for all
  languages and fonts, including non-Unicode Indic fonts (CDAC, Kruti Dev)
  that text-based importers cannot map back correctly.
"""
from __future__ import annotations

import io
import logging
import subprocess
import tempfile
from pathlib import Path

from app.utils import PDFProcessingError
from app.utils.concurrency import resolve_libreoffice_path
from app.utils.file_utils import sanitize_filename

logger = logging.getLogger(__name__)


class PptxService:
    @staticmethod
    def pptx_to_pdf(pptx_bytes: bytes, original_filename: str) -> bytes:
        return powerpoint_to_pdf(pptx_bytes, original_filename)

    @staticmethod
    def pdf_to_pptx(pdf_bytes: bytes, dpi: int = 150) -> bytes:
        return pdf_to_powerpoint(pdf_bytes, dpi=dpi)


def powerpoint_to_pdf(pptx_bytes: bytes, original_filename: str) -> bytes:
    """Convert a .ppt/.pptx file to PDF using LibreOffice headless."""
    libreoffice_path = resolve_libreoffice_path()
    if not libreoffice_path:
        raise PDFProcessingError("LibreOffice is not installed on the server.")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        input_path = tmp_path / sanitize_filename(original_filename)
        input_path.write_bytes(pptx_bytes)

        import uuid as _uuid
        profile_dir = tmp_path / f"lo_profile_{_uuid.uuid4().hex[:8]}"
        cmd = [
            libreoffice_path,
            "--headless",
            f"-env:UserInstallation=file:///{profile_dir.as_posix().lstrip('/')}",
            "--convert-to", "pdf",
            "--outdir", str(tmp_path),
            str(input_path),
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        except subprocess.TimeoutExpired:
            raise PDFProcessingError("PowerPoint to PDF conversion timed out.")

        if result.returncode != 0:
            logger.error("LibreOffice PPTX-to-PDF failed: %s", result.stderr[:500])
            raise PDFProcessingError(
                "PowerPoint to PDF conversion failed. Please try a different file."
            )

        pdfs = list(tmp_path.glob("*.pdf"))
        if not pdfs:
            raise PDFProcessingError("Conversion completed but no PDF was generated.")
        return pdfs[0].read_bytes()


def _pdf_to_pptx_libreoffice(pdf_bytes: bytes):
    """Try LibreOffice PDF→PPTX. Returns pptx bytes or None on failure."""
    libreoffice_path = resolve_libreoffice_path()
    if not libreoffice_path:
        return None

    try:
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            input_path = tmp_path / "input.pdf"
            input_path.write_bytes(pdf_bytes)

            import uuid as _uuid
            profile_dir = tmp_path / f"lo_profile_{_uuid.uuid4().hex[:8]}"
            cmd = [
                libreoffice_path,
                "--headless",
                f"-env:UserInstallation=file:///{profile_dir.as_posix().lstrip('/')}",
                "--infilter=impress_pdf_import",
                "--convert-to", "pptx",
                "--outdir", str(tmp_path),
                str(input_path),
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode != 0:
                logger.warning("LibreOffice PDF-to-PPTX failed: %s", result.stderr[:500])
                return None

            pptx_files = list(tmp_path.glob("*.pptx"))
            if not pptx_files:
                logger.warning("LibreOffice PDF-to-PPTX produced no output")
                return None

            return pptx_files[0].read_bytes()

    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice PDF-to-PPTX timed out")
        return None
    except Exception as exc:
        logger.warning("LibreOffice PDF-to-PPTX error: %s", exc)
        return None


def _pdf_to_pptx_images(pdf_bytes: bytes, dpi: int = 150) -> bytes:
    """Fallback: render each PDF page as an image slide."""
    try:
        from pdf2image import convert_from_bytes
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:
        raise PDFProcessingError(
            f"Missing dependency for PDF to PowerPoint: {e}. "
            "Install python-pptx and pdf2image."
        )

    from app.config import settings
    poppler_path = settings.POPPLER_PATH or None
    try:
        images = convert_from_bytes(pdf_bytes, dpi=dpi, poppler_path=poppler_path)
    except Exception as e:
        raise PDFProcessingError(f"Failed to render PDF pages: {e}")

    if not images:
        raise PDFProcessingError("PDF has no pages to convert.")

    prs = Presentation()

    first_w, first_h = images[0].size
    slide_h_inches = 7.5
    slide_w_inches = slide_h_inches * first_w / first_h
    prs.slide_width = Emu(int(slide_w_inches * 914400))
    prs.slide_height = Emu(int(slide_h_inches * 914400))

    blank_layout = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank_layout)

        buf = io.BytesIO()
        if img.mode == "RGBA":
            img = img.convert("RGB")
        img.save(buf, format="JPEG", quality=80, optimize=True)
        buf.seek(0)

        slide.shapes.add_picture(
            buf, Emu(0), Emu(0),
            width=prs.slide_width, height=prs.slide_height,
        )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def pdf_to_powerpoint(pdf_bytes: bytes, dpi: int = 150) -> bytes:
    """Convert PDF to PPTX using image-based rendering.

    Each PDF page is rendered as a high-quality image and placed on a
    slide whose dimensions match the original page aspect ratio.  This
    guarantees visual fidelity for all languages, fonts, and layouts —
    including non-Unicode fonts (e.g. CDAC_GISTSurekh for Hindi) that
    LibreOffice's PDF import cannot map back to text correctly.
    """
    return _pdf_to_pptx_images(pdf_bytes, dpi=dpi)
